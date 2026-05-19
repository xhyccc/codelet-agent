from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BLUE = RGBColor(0x1A, 0x3C, 0x6E)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox

# ── Slide 1: Title ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
add_textbox(slide, 1, 2.0, 11.3, 1.5, '网易新闻速览', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.8, 11.3, 1.0, 'Latest Headlines from 163.com', font_size=28, bold=False, color=RGBColor(0xBB,0xCC,0xDD), alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 5.2, 11.3, 0.6, '2026年5月18日', font_size=20, bold=False, color=RGBColor(0x99,0xAA,0xBB), alignment=PP_ALIGN.CENTER)

# ── Slide 2: Top Headlines ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, '📰 今日要闻', font_size=36, bold=True, color=DARK_BLUE)

headlines = [
    '中美元首北京会晤丨历史性的一页，要用历史的长镜头去端详',
    '特朗普回国已三天 \u201c不可思议的中国\u201d依然刷屏海外',
    '特朗普严厉警告\u201c台独\u201d后 赖清德首度发声被指态度强硬',
    '媒体：刚结束访华 特朗普打破惯例提\u201c四不\u201d示警\u201c台独\u201d',
    '在轨近200天 神二十一乘组\u201c太空出差\u201dVlog上新',
]

y = 1.5
for i, h in enumerate(headlines, 1):
    add_textbox(slide, 1.0, y, 0.6, 0.5, f'{i}.', font_size=22, bold=True, color=ACCENT_RED)
    add_textbox(slide, 1.6, y, 10.5, 0.5, h, font_size=20, bold=False, color=DARK_GRAY)
    y += 0.85

# ── Slide 3: Domestic News ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, '🇨🇳 国内新闻', font_size=36, bold=True, color=DARK_BLUE)

domestic = [
    '南水北调已累计向北方调水880亿立方米',
    '开局之年看中国丨向海而兴，向绿而行',
    '湖南：深化人工智能产业赋能和场景应用',
    '让收藏在博物馆里的文物活起来 中俄多彩主题年，拉紧人文交流纽带',
    '老太参团旅游拿出100万投资 子女报警后她又投了100万',
    '博主举报多名高校教授：他们求我别曝光 我说没办法',
    '幼童进手术室9小时死亡 法医刘良时隔20年再拿解剖刀',
]

y = 1.5
for i, h in enumerate(domestic, 1):
    add_textbox(slide, 1.0, y, 0.6, 0.5, f'{i}.', font_size=20, bold=True, color=ACCENT_RED)
    add_textbox(slide, 1.6, y, 10.5, 0.5, h, font_size=18, bold=False, color=DARK_GRAY)
    y += 0.75

# ── Slide 4: International & Business ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, '🌍 国际 & 财经', font_size=36, bold=True, color=DARK_BLUE)

intl_biz = [
    '69年来首次亏损4000亿日元，本田电动化猛踩刹车',
    '英媒：中国两轮车企借绿色浪潮驶入欧洲',
    '日本零食包装也变黑白色 高市早苗还在强行\u201c挽尊\u201d',
    '国产跑鞋里的\u201c中国服务\u201d硬实力',
    '狄莺儿子被捕:母子同床15年 儿子有生理反应才分床',
]

y = 1.5
for i, h in enumerate(intl_biz, 1):
    add_textbox(slide, 1.0, y, 0.6, 0.5, f'{i}.', font_size=20, bold=True, color=ACCENT_RED)
    add_textbox(slide, 1.6, y, 10.5, 0.5, h, font_size=18, bold=False, color=DARK_GRAY)
    y += 0.85

# ── Slide 5: Tech & Culture ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_textbox(slide, 0.8, 0.4, 11.7, 0.8, '🔬 科技 & 文化', font_size=36, bold=True, color=DARK_BLUE)

tech_culture = [
    '国产大模型集体更新后能力有多强？',
    '湖南：深化人工智能产业赋能和场景应用',
    '我们为什么需要博物馆',
    '【光明论坛】在\u201c赶大集\u201d中触摸乡村消费脉搏',
    '5招纠正你的错误体态，低头族提升气质必备',
]

y = 1.5
for i, h in enumerate(tech_culture, 1):
    add_textbox(slide, 1.0, y, 0.6, 0.5, f'{i}.', font_size=20, bold=True, color=ACCENT_RED)
    add_textbox(slide, 1.6, y, 10.5, 0.5, h, font_size=18, bold=False, color=DARK_GRAY)
    y += 0.85

# ── Slide 6: Closing ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_textbox(slide, 1, 2.5, 11.3, 1.2, '谢谢观看', font_size=48, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 4.0, 11.3, 0.8, '数据来源：网易新闻 news.163.com', font_size=20, bold=False, color=RGBColor(0x99,0xAA,0xBB), alignment=PP_ALIGN.CENTER)

output_path = '网易新闻速览_20260518.pptx'
prs.save(output_path)
print(f'Saved to {output_path}')
